from pptx import Presentation
path="C:/Users/MECHREVO/Desktop/《窗外，真的在下雨吗？.pptx"
ppt=Presentation(r"C:/Users/MECHREVO/Desktop/《窗外，真的在下雨吗？.pptx")
layout=ppt.slide_layouts[0]
slideall=ppt.slides
newslide=slideall.add_slide(layout)#在所有幻灯片末尾加入并输入版式
for single_plh in newslide.placeholders:#获得占位符对象
    phf=single_plh.placeholder_format#得到对象的属性
    print(f"{phf.idx}-{phf.type}")#idx为编码 type为类型
allplh=newslide.placeholders
allplh[0].text="我爱你"#索引加入文本 与word加入文本不太一样
ppt.save("C:/Users/MECHREVO/Desktop/《窗外，真的在下雨吗？.pptx")
#导入文件名 import os
#items=os.listdir(文件夹)
#for item in items:
#   filename=itme.split(".")[0]
#   print(filename)