from pptx import Presentation
import docx
file=docx.Document()
path=r"C:/Users/MECHREVO\Desktop/《窗外，真的在下雨吗？.pptx"
pptxfile=Presentation(path)
n=1
for slide in pptxfile.slides: 
     file.add_heading(f"第{n}页",level=2)
     n=n+1
     for shape in slide.shapes :
          if shape.has_text_frame==True:
             textframe=shape.text_frame
             for paragraph in textframe.paragraphs:
                 for run in paragraph.runs:
                     texts=run.text
                     file.add_paragraph(texts)
file.save("E:/python/2.docx")
